import asyncio
import hashlib
import json
import logging
import os
import re
import tempfile
import zipfile
from collections import defaultdict
from datetime import UTC, datetime
from pathlib import Path, PurePosixPath
from typing import Annotated
from urllib.parse import quote

from fastapi import APIRouter, Depends, HTTPException, Path as PathParam, Query, Request

from app.services import event_hooks
from fastapi.responses import FileResponse as FastAPIFileResponse
from fastapi.responses import Response, StreamingResponse
from sqlalchemy import and_, func, or_
from sqlalchemy.orm import Session

import app.config as config
from app.auth import check_drive_access, get_unlocked_groups
from app.database import get_db
from app.models import (
    File,
    FileExif,
    FileRelation,
    Tag,
    active_file_filter,
    file_tags,
)
from app.schemas import (
    ArchiveContentsResponse,
    ArchiveEntryResponse,
    BatchCopyRequest,
    BatchCopyResponse,
    BatchIdsRequest,
    BatchMoveRequest,
    BatchPurgeResponse,
    BatchRenameRequest,
    BatchRenameResponse,
    BatchRestoreResponse,
    BatchTagRequest,
    ExifResponse,
    FileCopyRequest,
    FileRelationItem,
    FileRelationsResponse,
    FileResponse,
    FileMoveRequest,
    FileRenameRequest,
    FileUpdate,
    NeighborsResponse,
    RelatedFileSummary,
    SubtitleInfo,
    TagUpdate,
    file_to_response,
)
from app.services import fileops
from app.services import ws as ws_service
from app.services.filetype import classify
from app.services.frontmatter import (
    compose as compose_frontmatter,
    ensure_id,
    extract_valid_aliases,
    extract_valid_tags,
    parse as parse_frontmatter,
)
from app.services.markdown_relations import (
    ResolveDiagnostic,
    extract_links,
    resolve_wiki_targets,
    resolve_wiki_targets_with_map,
)
from app.services.heic import HEIC_MIME_TYPES, convert_heic_to_jpeg
from app.services.subtitle import convert_srt_to_vtt, detect_subtitles

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/files", tags=["files"])

FileId = Annotated[str, PathParam(min_length=12, max_length=12, pattern=r"^[A-Za-z0-9_-]+$")]

PLACEHOLDER_THUMBNAIL = Path(__file__).parent.parent / "static" / "placeholder.jpg"

_ARCHIVE_ENTRY_MAX_SIZE = 50 * 1024 * 1024  # 50MB
_MAX_ARCHIVE_ENTRIES = 10_000
_archive_semaphore = asyncio.Semaphore(3)

_TEXT_WRITE_ALLOWED_MIMES = frozenset({"text/markdown", "text/plain"})
_TEXT_WRITE_MAX_BYTES = 1 * 1024 * 1024  # 1 MB
_text_write_locks: dict[str, asyncio.Lock] = defaultdict(asyncio.Lock)

_DANGEROUS_INLINE_MIMES = frozenset({
    "text/html",
    "application/xhtml+xml",
    "image/svg+xml",
    "text/xml",
    "application/xml",
    "application/xslt+xml",
})

def _sync_md_file_relations(
    db: Session,
    file_id: str,
    drive: str,
    content: str,
    self_dir: str,
) -> list[ResolveDiagnostic]:
    """Sync file_relations for a .md file using its body and frontmatter.

    Spec ``2026-05-12-markdown-link-three-forms.md`` §3.5.

    Three sources contribute to ``kind='related'`` relations on ``.md`` files:

    * ``loft://<id>`` in body — direct file id reference.
    * ``[[X]]`` in body — wiki target, resolved drive-scoped via
      :func:`resolve_wiki_targets` (see §3.3 for precedence).
    * ``source_file_ids: [...]`` in frontmatter — raw file IDs (same format
      as distill / Ask-save notes). Allows relations to persist even when the
      body is empty or the user hasn't written explicit links yet.

    Any combination of these sources can maintain a relation independently;
    removing one source only removes the relation if no other source references
    the same file.

    Ambiguous and unresolved wiki targets do **not** create relations
    (per spec §7.5 — auto-picking would surprise the writer). They are
    returned as diagnostics so the caller can surface a warning.

    Same-drive + active filter is applied; cross-drive references are
    silently dropped (drive = security boundary).
    """
    extracted = extract_links(content)
    loft_ids = {m for m in extracted.loft_ids if m != file_id}
    wiki_ids, diagnostics = resolve_wiki_targets(
        db, drive, self_dir, extracted.wiki_targets
    )

    # frontmatter source_file_ids contribute the same way as loft_ids.
    # Malformed frontmatter is silently ignored — a save must never fail
    # due to a YAML parse error in this projection step.
    fm_ids: set[str] = set()
    try:
        _parsed_fm = parse_frontmatter(content)
        _raw = _parsed_fm.metadata.get("source_file_ids")
        if isinstance(_raw, list):
            fm_ids = {
                str(v) for v in _raw
                if isinstance(v, str) and v and v != file_id
            }
    except Exception:
        pass

    # loft_ids and fm_ids bypass the resolver, so apply the same-drive +
    # active filter explicitly. Without this, a stale id pointing at a
    # different drive could leak into ``file_relations``.
    raw_ids = (loft_ids | fm_ids) - {file_id}
    valid_raw: set[str] = set()
    if raw_ids:
        rows = (
            db.query(File.id)
            .filter(
                File.id.in_(raw_ids),
                File.drive == drive,
                active_file_filter(),
            )
            .all()
        )
        valid_raw = {row.id for row in rows}

    target_ids = (valid_raw | wiki_ids) - {file_id}

    # Current file_relations for this file (kind='related', both directions)
    existing_rels = db.query(FileRelation).filter(
        or_(
            FileRelation.file_id_a == file_id,
            FileRelation.file_id_b == file_id,
        ),
        FileRelation.kind == "related",
    ).all()
    existing_map: dict[str, FileRelation] = {}
    for rel in existing_rels:
        other = rel.file_id_b if rel.file_id_a == file_id else rel.file_id_a
        existing_map[other] = rel

    existing_ids = set(existing_map)
    to_add = target_ids - existing_ids
    to_remove = existing_ids - target_ids

    for tid in to_add:
        db.add(
            FileRelation(
                file_id_a=file_id,
                file_id_b=tid,
                kind="related",
                created_at=datetime.now(UTC),
            )
        )
    for tid in to_remove:
        db.delete(existing_map[tid])

    return diagnostics


def _is_markdown_file(file: File) -> bool:
    """Whether frontmatter should be parsed on content writes for ``file``.

    Mirrors the frontend ``isMarkdown`` heuristic (frontend/src/lib/tags.ts):
    trust ``text/markdown`` first, fall back to the ``.md`` extension
    because some older rows still report ``text/plain`` for ``.md``.
    Keeping the two sides aligned is a spec §D1 requirement — a file
    that the UI treats as markdown must project frontmatter on the
    backend, and vice versa.
    """
    if (file.mime_type or "") == "text/markdown":
        return True
    return file.filename.lower().endswith(".md")


def replace_file_tags(db: Session, file: File, tag_names: list[str]) -> None:
    """Replace ``file.tags`` with the given names, reusing existing Tag rows.

    Shared by ``PUT /api/files/{id}/tags``, ``PUT /api/files/batch/tags`` and
    the internal ``POST /api/internal/files/{id}/tags`` (spec
    ``2026-04-24-knowledge-tag-unification.md``). Case-insensitive dedup
    via ``func.lower(Tag.name)``; the ``Tag`` namespace is per-drive
    (``uq_tags_drive_name``).

    SECURITY: callers MUST verify drive access before invoking this
    helper. It performs no authorisation check — it trusts that
    ``_get_file_or_404`` or equivalent was called upstream.

    Transactional contract: the caller is responsible for ``db.commit()``
    and for invoking ``cleanup_orphan_tags`` afterwards. Kept
    commit-free so batch callers can replace tags on many files in a
    single transaction.
    """
    tag_objects: list[Tag] = []
    for tag_name in tag_names:
        tag = (
            db.query(Tag)
            .filter(
                func.lower(Tag.name) == tag_name.lower(),
                Tag.drive == file.drive,
            )
            .first()
        )
        if not tag:
            tag = Tag(name=tag_name, drive=file.drive)
        elif tag.name != tag_name:
            tag.name = tag_name
            db.add(tag)
            db.flush()
        tag_objects.append(tag)
    file.tags = tag_objects


def cleanup_orphan_tags(db: Session) -> int:
    """Remove Tag rows no longer referenced by any file. Returns count deleted.

    Transactional contract: caller commits. Symmetric with
    ``replace_file_tags`` so the two helpers always compose inside a
    single transaction.

    ``db.flush()`` first so pending ``file.tags = [...]`` reassignments
    from ``replace_file_tags`` are written to the ``file_tags`` table
    before the OUTER JOIN query runs. Without the flush, SQLAlchemy
    keeps the association change in the session and the orphan query
    reads a stale snapshot.
    """
    db.flush()
    orphans = (
        db.query(Tag)
        .outerjoin(file_tags)
        .filter(file_tags.c.file_id.is_(None))
        .all()
    )
    for orphan in orphans:
        db.delete(orphan)
    return len(orphans)




def _decode_zip_filename(info: zipfile.ZipInfo) -> str:
    """Decode ZIP entry filename, handling Shift_JIS encoded names.

    ZIP files created on Japanese Windows encode filenames in Shift_JIS (CP932)
    but don't set the UTF-8 flag. Python's zipfile decodes them as CP437,
    producing garbled text. This function detects and re-decodes as CP932.
    """
    # If UTF-8 flag is set, Python already decoded correctly
    if info.flag_bits & 0x800:
        return info.filename

    # Try re-encoding from CP437 back to bytes, then decode as CP932
    try:
        raw = info.filename.encode("cp437")
    except UnicodeEncodeError:
        return info.filename

    # Pure ASCII is identical in both encodings — no need to re-decode
    if all(b < 0x80 for b in raw):
        return info.filename

    try:
        return raw.decode("cp932")
    except UnicodeDecodeError:
        # Not Shift_JIS — return as-is (original CP437 decode)
        return info.filename


_SAFE_INLINE_TYPES = frozenset({
    "image/jpeg", "image/png", "image/gif", "image/webp", "image/avif",
    "image/bmp", "text/plain",
})


def _validate_path(file_path: str, base_dir: Path) -> Path:
    real_path = Path(os.path.realpath(file_path))
    real_base = Path(os.path.realpath(base_dir))
    base_str = str(real_base)
    if not (str(real_path) == base_str or str(real_path).startswith(base_str + os.sep)):
        raise HTTPException(status_code=403, detail="Access denied")
    return real_path


_to_response = file_to_response


def _detect_file_subtitles(file: File) -> list[SubtitleInfo]:
    if file.file_type != "video" and file.mime_type != "application/vnd.litloft.loft+json":
        return []
    drive_path = config.get_drive_path(file.drive)
    raw = detect_subtitles(file.file_path, drive_path)
    return [
        SubtitleInfo(index=i, language=s["language"], format=s["format"], label=s["label"])
        for i, s in enumerate(raw)
    ]


def _is_drive_accessible(drive_name: str, unlocked_groups: list[str]) -> bool:
    access_group = config.get_drive_access_group(drive_name)
    return not access_group or access_group in unlocked_groups


def _get_file_or_404(
    db: Session, file_id: str, unlocked_groups: list[str]
) -> File:
    file = db.query(File).filter(File.id == file_id, active_file_filter()).first()
    if not file:
        raise HTTPException(status_code=404, detail="File not found")
    check_drive_access(file.drive, unlocked_groups)
    return file


def _get_trashed_file_or_404(
    db: Session, file_id: str, unlocked_groups: list[str]
) -> File:
    file = db.query(File).filter(File.id == file_id, File.deleted_at.isnot(None)).first()
    if not file:
        raise HTTPException(status_code=404, detail="File not found in trash")
    check_drive_access(file.drive, unlocked_groups)
    return file


def _get_missing_file_or_404(
    db: Session, file_id: str, unlocked_groups: list[str]
) -> File:
    file = db.query(File).filter(
        File.id == file_id,
        File.missing_since.isnot(None),
        File.deleted_at.is_(None),
    ).first()
    if not file:
        raise HTTPException(status_code=404, detail="File not found in missing")
    check_drive_access(file.drive, unlocked_groups)
    return file


def _get_trashed_or_missing_file_or_404(
    db: Session, file_id: str, unlocked_groups: list[str]
) -> File:
    """Get a file that is either trashed or missing (both purge-eligible)."""
    file = db.query(File).filter(
        File.id == file_id,
        or_(File.deleted_at.isnot(None), File.missing_since.isnot(None)),
    ).first()
    if not file:
        raise HTTPException(status_code=404, detail="File not found")
    check_drive_access(file.drive, unlocked_groups)
    return file


def _get_file_any_state_or_404(
    db: Session, file_id: str, unlocked_groups: list[str]
) -> File:
    """Get file regardless of state (for thumbnail/metadata access in trash or missing)."""
    file = db.query(File).filter(File.id == file_id).first()
    if not file:
        raise HTTPException(status_code=404, detail="File not found")
    check_drive_access(file.drive, unlocked_groups)
    return file


@router.post("/batch/get", response_model=list[FileResponse])
async def batch_get(
    body: BatchIdsRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    files = db.query(File).filter(File.id.in_(body.ids), active_file_filter()).all()
    file_map = {f.id: f for f in files}
    return [
        _to_response(file_map[fid])
        for fid in body.ids
        if fid in file_map and _is_drive_accessible(file_map[fid].drive, unlocked_groups)
    ]


@router.post("/batch/delete")
async def batch_delete(
    body: BatchIdsRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    deleted = 0
    deleted_ids = []
    errors = []
    for file_id in body.ids:
        try:
            _get_file_or_404(db, file_id, unlocked_groups)
            fileops.delete_file(db, file_id)
            deleted += 1
            deleted_ids.append(file_id)
        except HTTPException as e:
            errors.append({"id": file_id, "error": e.detail})
    if deleted_ids:
        asyncio.create_task(
            event_hooks.emit("files.deleted", {"file_ids": deleted_ids, "type": "soft_delete"})
        )
    return {"deleted": deleted, "errors": errors}


@router.put("/batch/move")
async def batch_move(
    body: BatchMoveRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    moved = 0
    moved_ids: list[str] = []
    errors = []
    try:
        for file_id in body.ids:
            try:
                _get_file_or_404(db, file_id, unlocked_groups)
                fileops.move_file(db, file_id, body.target_drive, body.target_folder_path)
                moved += 1
                moved_ids.append(file_id)
            except HTTPException as e:
                errors.append({"id": file_id, "error": e.detail})
    finally:
        # Emit even if an unexpected exception aborts the loop: per-file
        # ``move_file`` commits individually, so ids already in ``moved_ids``
        # are durable on disk and need to reach addons regardless.
        if moved_ids:
            await event_hooks.emit("files.moved", {"file_ids": moved_ids})
    return {"moved": moved, "errors": errors}


@router.put("/batch/tags")
async def batch_tags(
    body: BatchTagRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    updated = 0
    updated_ids = []
    errors = []
    for file_id in body.ids:
        try:
            file = _get_file_or_404(db, file_id, unlocked_groups)
            replace_file_tags(db, file, body.tags)
            updated += 1
            updated_ids.append(file_id)
        except HTTPException as e:
            errors.append({"id": file_id, "error": e.detail})
    cleanup_orphan_tags(db)
    db.commit()
    if updated_ids:
        asyncio.create_task(
            event_hooks.emit("files.updated", {"file_ids": updated_ids})
        )
    return {"updated": updated, "errors": errors}


@router.put("/batch/rename", response_model=BatchRenameResponse)
async def batch_rename(
    body: BatchRenameRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    files = []
    for file_id in body.ids:
        files.append(_get_file_or_404(db, file_id, unlocked_groups))

    kwargs = body.model_dump(
        exclude={"ids", "mode"},
        exclude_none=True,
    )
    results = fileops.batch_rename(db, files, body.mode, **kwargs)
    renamed_ids = [
        r["id"] for r in results if r.get("old_name") != r.get("new_name")
    ]
    if renamed_ids:
        await event_hooks.emit("files.moved", {"file_ids": renamed_ids})
    return {"renamed": len(results), "results": results}


@router.post("/batch/restore", response_model=BatchRestoreResponse)
async def batch_restore(
    body: BatchIdsRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    restored = 0
    restored_ids = []
    errors = []
    for file_id in body.ids:
        try:
            _get_trashed_file_or_404(db, file_id, unlocked_groups)
            fileops.restore_file(db, file_id)
            restored += 1
            restored_ids.append(file_id)
        except HTTPException as e:
            errors.append({"id": file_id, "error": e.detail})
    if restored_ids:
        asyncio.create_task(
            event_hooks.emit("files.restored", {"file_ids": restored_ids})
        )
    return {"restored": restored, "errors": errors}


@router.post("/batch/purge", response_model=BatchPurgeResponse)
async def batch_purge(
    body: BatchIdsRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    purged = 0
    purged_ids = []
    errors = []
    for file_id in body.ids:
        try:
            file = _get_trashed_or_missing_file_or_404(db, file_id, unlocked_groups)
            if file.missing_since is not None and file.deleted_at is None:
                fileops.purge_missing_file(db, file_id)
            else:
                fileops.purge_file(db, file_id)
            purged += 1
            purged_ids.append(file_id)
        except HTTPException as e:
            errors.append({"id": file_id, "error": e.detail})
    if purged_ids:
        asyncio.create_task(
            event_hooks.emit("files.purged", {"file_ids": purged_ids})
        )
    return {"purged": purged, "errors": errors}


@router.post("/batch/copy", response_model=BatchCopyResponse)
async def batch_copy(
    body: BatchCopyRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    copied = 0
    copied_ids = []
    errors = []
    for file_id in body.ids:
        try:
            _get_file_or_404(db, file_id, unlocked_groups)
            new_file = fileops.copy_file(db, file_id, body.target_drive, body.target_folder_path)
            copied_ids.append(new_file.id)
            copied += 1
        except HTTPException as e:
            errors.append({"id": file_id, "error": e.detail})
    if copied_ids:
        asyncio.create_task(
            event_hooks.emit("files.created", {"file_ids": copied_ids})
        )
    return {"copied": copied, "errors": errors}


@router.get("/{file_id}", response_model=FileResponse)
async def get_file(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)
    subtitles = _detect_file_subtitles(file)
    return _to_response(file, subtitles=subtitles)


@router.get("/{file_id}/neighbors", response_model=NeighborsResponse)
async def get_file_neighbors(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
    sort: str = Query("created_at", pattern="^(created_at|title|file_size|likes)$"),
    order: str = Query("desc", pattern="^(asc|desc)$"),
):
    file = _get_file_or_404(db, file_id, unlocked_groups)

    sort_col = getattr(File, sort)
    current_val = getattr(file, sort)

    base = db.query(File.id).filter(
        File.drive == file.drive,
        File.folder_path == file.folder_path,
        File.id != file.id,
        active_file_filter(),
    )

    if order == "asc":
        prev_query = base.filter(
            or_(
                sort_col < current_val,
                and_(sort_col == current_val, File.id < file.id),
            )
        ).order_by(sort_col.desc(), File.id.desc()).limit(1)

        next_query = base.filter(
            or_(
                sort_col > current_val,
                and_(sort_col == current_val, File.id > file.id),
            )
        ).order_by(sort_col.asc(), File.id.asc()).limit(1)
    else:
        prev_query = base.filter(
            or_(
                sort_col > current_val,
                and_(sort_col == current_val, File.id > file.id),
            )
        ).order_by(sort_col.asc(), File.id.asc()).limit(1)

        next_query = base.filter(
            or_(
                sort_col < current_val,
                and_(sort_col == current_val, File.id < file.id),
            )
        ).order_by(sort_col.desc(), File.id.desc()).limit(1)

    prev_row = prev_query.first()
    next_row = next_query.first()

    return NeighborsResponse(
        prev_id=prev_row[0] if prev_row else None,
        next_id=next_row[0] if next_row else None,
    )


@router.put("/{file_id}", response_model=FileResponse)
async def update_file(
    file_id: FileId,
    update: FileUpdate,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)

    update_data = update.model_dump(exclude_unset=True)
    for key, value in update_data.items():
        setattr(file, key, value)

    db.commit()
    db.refresh(file)
    asyncio.create_task(
        event_hooks.emit("files.updated", {"file_ids": [file_id]})
    )
    return _to_response(file)


@router.post("/{file_id}/like", response_model=FileResponse)
async def like_file(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)

    file.likes = File.likes + 1
    db.commit()
    db.refresh(file)
    return _to_response(file)


@router.post("/{file_id}/dislike", response_model=FileResponse)
async def dislike_file(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)

    file.likes = File.likes - 1
    db.commit()
    db.refresh(file)
    return _to_response(file)


@router.post("/{file_id}/favorite", response_model=FileResponse)
async def toggle_favorite(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)

    file.is_favorite = not file.is_favorite
    db.commit()
    db.refresh(file)
    asyncio.create_task(
        event_hooks.emit("files.updated", {"file_ids": [file_id]})
    )
    return _to_response(file)


@router.put("/{file_id}/tags", response_model=FileResponse)
async def update_file_tags(
    file_id: FileId,
    update: TagUpdate,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)
    replace_file_tags(db, file, update.tags)
    cleanup_orphan_tags(db)
    db.commit()
    db.refresh(file)
    asyncio.create_task(
        event_hooks.emit("files.updated", {"file_ids": [file_id]})
    )
    return _to_response(file)


@router.get("/{file_id}/stream")
async def stream_file(
    file_id: FileId,
    request: Request,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
    download: bool = False,
):
    file = _get_file_any_state_or_404(db, file_id, unlocked_groups)
    if file.missing_since is not None and file.deleted_at is None:
        raise HTTPException(status_code=410, detail="File is missing")

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(
        str(drive_path / file.file_path), drive_path
    )
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")

    content_type = file.mime_type or "application/octet-stream"

    # HEIC/HEIF: serve converted JPEG for browser compatibility
    if content_type in HEIC_MIME_TYPES:
        jpeg_path = convert_heic_to_jpeg(str(file_path), config.CONVERTED_DIR)
        if jpeg_path is not None:
            file_path = jpeg_path
            content_type = "image/jpeg"

    file_size = file_path.stat().st_size
    range_header = request.headers.get("range")

    if range_header:
        try:
            range_spec = range_header.replace("bytes=", "").split(",")[0]
            parts = range_spec.split("-")
            start = int(parts[0])
            end = (
                int(parts[1])
                if parts[1]
                else min(start + config.CHUNK_SIZE - 1, file_size - 1)
            )
            end = min(end, file_size - 1)
            if start < 0 or start > end or start >= file_size:
                raise HTTPException(status_code=416, detail="Range not satisfiable")
        except HTTPException:
            raise
        except (ValueError, IndexError):
            raise HTTPException(status_code=416, detail="Invalid range header")

        def iter_chunks():
            with open(file_path, "rb") as f:
                f.seek(start)
                remaining = end - start + 1
                while remaining > 0:
                    chunk = f.read(min(config.CHUNK_SIZE, remaining))
                    if not chunk:
                        break
                    remaining -= len(chunk)
                    yield chunk

        headers = {
            "Content-Range": f"bytes {start}-{end}/{file_size}",
            "Accept-Ranges": "bytes",
            "Content-Length": str(end - start + 1),
            "Content-Type": content_type,
            "X-Content-Type-Options": "nosniff",
        }
        if download or content_type in _DANGEROUS_INLINE_MIMES:
            headers["Content-Disposition"] = f"attachment; filename*=UTF-8''{quote(file.filename, safe='')}"

        return StreamingResponse(iter_chunks(), status_code=206, headers=headers)

    # Small text files: serve as full body with ETag so clients can use
    # the content-hash for optimistic locking on PUT /content without having
    # to hash on the client (crypto.subtle is unavailable in non-secure contexts
    # like HTTP over LAN IPs).
    if (
        (file.mime_type or "") in _TEXT_WRITE_ALLOWED_MIMES
        and file_size <= _TEXT_WRITE_MAX_BYTES
    ):
        data = file_path.read_bytes()
        headers = {
            "Accept-Ranges": "bytes",
            "Content-Length": str(len(data)),
            "Content-Type": content_type,
            "ETag": f'"{_compute_text_etag(data)}"',
            "X-Content-Type-Options": "nosniff",
        }
        if download or content_type in _DANGEROUS_INLINE_MIMES:
            headers["Content-Disposition"] = f"attachment; filename*=UTF-8''{quote(file.filename, safe='')}"
        return Response(content=data, headers=headers)

    def iter_full():
        with open(file_path, "rb") as f:
            while chunk := f.read(config.CHUNK_SIZE):
                yield chunk

    headers = {
        "Accept-Ranges": "bytes",
        "Content-Length": str(file_size),
        "Content-Type": content_type,
        "X-Content-Type-Options": "nosniff",
    }
    if download or content_type in _DANGEROUS_INLINE_MIMES:
        headers["Content-Disposition"] = f"attachment; filename*=UTF-8''{quote(file.filename, safe='')}"

    return StreamingResponse(iter_full(), headers=headers)


_RENDER_MAX_BYTES = 5 * 1024 * 1024  # 5 MB

# Inline script appended to HTML responses so the iframe can report its
# rendered height to the parent window via postMessage. The parent listens
# for `litloft:height` messages and resizes the iframe accordingly, which
# yields a single-scroll layout (the page scrolls, the iframe does not).
# Skipped in fullscreen mode (#litloft-fullscreen hash), where the iframe
# already occupies the viewport.
_RENDER_BOOTSTRAP = (
    b"<script>(function(){"
    b"if(location.hash==='#litloft-fullscreen')return;"
    b"function s(){parent.postMessage({type:'litloft:height',"
    b"value:Math.ceil(document.documentElement.scrollHeight)},'*');}"
    b"new ResizeObserver(s).observe(document.documentElement);s();"
    b"})();</script>"
)

# CSP for /render responses. The `sandbox` directive forces a null origin
# even on top-level navigation, so the document cannot read parent cookies
# or storage. The token list mirrors the iframe `sandbox` attribute:
# without `allow-scripts` the document cannot run any JS, so both the
# bootstrap resize script and the AI artifact's own code would silently
# fail. `allow-popups` keeps `<a target="_blank">` working. We intentionally
# omit `allow-same-origin`, `allow-top-navigation`, and
# `allow-popups-to-escape-sandbox` to keep the document in an opaque origin.
# `default-src 'none'` blocks all network egress (fetch, WebSocket, beacon,
# EventSource); the explicit allowlists let common AI artifact CDNs load
# scripts/styles/fonts without opening generic exfil paths.
_RENDER_CSP = (
    "sandbox allow-scripts allow-popups; "
    "default-src 'none'; "
    "script-src 'unsafe-inline' 'unsafe-eval' "
    "https://cdn.jsdelivr.net https://unpkg.com https://esm.sh "
    "https://cdnjs.cloudflare.com https://cdn.tailwindcss.com; "
    "style-src 'unsafe-inline' https://cdn.jsdelivr.net https://fonts.googleapis.com; "
    "img-src 'self' data: blob:; "
    "font-src https://fonts.gstatic.com data:; "
    "connect-src 'none'; "
    "form-action 'none'; "
    "frame-ancestors 'self'"
)

# Case-insensitive locator for </body>. Falls back to end-of-document
# append when absent. The regex consumes the literal </body> so we can
# splice the bootstrap immediately before it.
_BODY_CLOSE_RE = re.compile(rb"</body\s*>", re.IGNORECASE)


def _inject_render_bootstrap(html_bytes: bytes) -> bytes:
    """Inject the height-reporting bootstrap script into an HTML response.

    Inserts before the final </body> when present (case-insensitive),
    otherwise appends to the end. Returns the modified bytes; callers
    must ensure the input decodes as UTF-8 before reaching this helper.
    """
    match = None
    for m in _BODY_CLOSE_RE.finditer(html_bytes):
        match = m
    if match is None:
        return html_bytes + _RENDER_BOOTSTRAP
    return html_bytes[: match.start()] + _RENDER_BOOTSTRAP + html_bytes[match.start() :]


@router.get("/{file_id}/render")
async def render_file(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    """Serve text/html in a sandboxed iframe with strict CSP.

    Companion to /stream which forces `Content-Disposition: attachment`
    for HTML to block top-level XSS via link luring. /render is a
    separate path intended only for the in-app file detail iframe;
    even on direct navigation, the CSP `sandbox` directive forces a
    null origin so parent cookies and storage stay isolated.
    """
    file = _get_file_any_state_or_404(db, file_id, unlocked_groups)
    if file.missing_since is not None and file.deleted_at is None:
        raise HTTPException(status_code=410, detail="File is missing")
    if (file.mime_type or "") != "text/html":
        raise HTTPException(status_code=404, detail="Not an HTML file")

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(str(drive_path / file.file_path), drive_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")

    file_size = file_path.stat().st_size
    if file_size > _RENDER_MAX_BYTES:
        raise HTTPException(status_code=413, detail="HTML file too large to render")

    raw = file_path.read_bytes()
    try:
        raw.decode("utf-8")
    except UnicodeDecodeError:
        raise HTTPException(
            status_code=415,
            detail="Only UTF-8 encoded HTML is supported",
        )

    body = _inject_render_bootstrap(raw)
    headers = {
        "Content-Type": "text/html; charset=utf-8",
        "Content-Disposition": "inline",
        "X-Content-Type-Options": "nosniff",
        "Content-Security-Policy": _RENDER_CSP,
        "Cache-Control": "no-store",
    }
    return Response(content=body, headers=headers)


_OFFICE_MIME_EXTENSIONS = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}
_PREVIEW_TEXT_MAX_CHARS = 400


def _extract_office_preview(file_path: Path, mime_type: str) -> str:
    ext = _OFFICE_MIME_EXTENSIONS.get(mime_type, "")
    try:
        if ext == ".docx":
            import docx
            doc = docx.Document(str(file_path))
            parts: list[str] = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)
                if sum(len(p) for p in parts) >= _PREVIEW_TEXT_MAX_CHARS:
                    break
            return "\n".join(parts)[:_PREVIEW_TEXT_MAX_CHARS]

        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            lines: list[str] = []
            ws = wb.worksheets[0] if wb.worksheets else None
            if ws:
                for row in ws.iter_rows(values_only=True, max_row=20):
                    cells = [str(c) for c in row if c is not None and str(c).strip()]
                    if cells:
                        lines.append(" ".join(cells))
                    if sum(len(l) for l in lines) >= _PREVIEW_TEXT_MAX_CHARS:
                        break
            wb.close()
            return "\n".join(lines)[:_PREVIEW_TEXT_MAX_CHARS]

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                parts.append(text)
                if sum(len(p) for p in parts) >= _PREVIEW_TEXT_MAX_CHARS:
                    break
            return "\n".join(parts)[:_PREVIEW_TEXT_MAX_CHARS]
    except Exception as e:
        logger.warning("Office preview extraction failed for %s: %s", file_path, e)
    return ""


@router.get("/{file_id}/preview-text")
async def get_preview_text(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_any_state_or_404(db, file_id, unlocked_groups)
    if file.mime_type not in _OFFICE_MIME_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Not an Office document")

    if file.missing_since is not None and file.deleted_at is None:
        raise HTTPException(status_code=410, detail="File is missing")

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(str(drive_path / file.file_path), drive_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")

    text = _extract_office_preview(file_path, file.mime_type)
    return Response(content=text, media_type="text/plain; charset=utf-8")


@router.get("/{file_id}/thumbnail")
async def get_thumbnail(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_any_state_or_404(db, file_id, unlocked_groups)

    if file.thumbnail_path:
        thumb_path = _validate_path(
            str(config.THUMBNAILS_DIR / file.thumbnail_path), config.DATA_DIR
        )
        if thumb_path.exists():
            return FastAPIFileResponse(str(thumb_path), media_type="image/jpeg")

    if PLACEHOLDER_THUMBNAIL.exists():
        return FastAPIFileResponse(
            str(PLACEHOLDER_THUMBNAIL),
            media_type="image/jpeg",
            headers={"Cache-Control": "no-cache"},
        )

    raise HTTPException(status_code=404, detail="Thumbnail not found")



@router.get("/{file_id}/archive", response_model=ArchiveContentsResponse)
async def get_archive_contents(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)
    if file.file_type != "archive":
        raise HTTPException(status_code=404, detail="File is not an archive")

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(str(drive_path / file.file_path), drive_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")

    entries = []
    total_size = 0
    with zipfile.ZipFile(str(file_path), "r") as zf:
        for info in zf.infolist():
            if len(entries) >= _MAX_ARCHIVE_ENTRIES:
                break
            # Skip symlink entries
            if info.external_attr != 0:
                mode = info.external_attr >> 16
                if mode != 0 and (mode & 0o170000) == 0o120000:
                    continue
            decoded_name = _decode_zip_filename(info)
            is_dir = info.is_dir()
            clean_path = decoded_name.rstrip("/") if is_dir else decoded_name
            entry_name = PurePosixPath(clean_path).name
            file_type, mime_type = classify(decoded_name) if not is_dir else ("other", "")
            entries.append(ArchiveEntryResponse(
                path=decoded_name,
                filename=entry_name,
                file_size=info.file_size,
                compressed_size=info.compress_size,
                file_type=file_type if not is_dir else "directory",
                mime_type=mime_type,
                is_dir=is_dir,
            ))
            total_size += info.file_size

    entries_sorted = sorted(entries, key=lambda e: e.path)
    return ArchiveContentsResponse(
        entries=entries_sorted,
        total_entries=len(entries_sorted),
        total_size=total_size,
    )


@router.get("/{file_id}/archive/entry")
async def get_archive_entry(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
    path: str = Query(..., min_length=1),
):
    clean = PurePosixPath(path)
    if clean.is_absolute() or ".." in clean.parts:
        raise HTTPException(status_code=400, detail="Invalid entry path")

    file = _get_file_or_404(db, file_id, unlocked_groups)
    if file.file_type != "archive":
        raise HTTPException(status_code=404, detail="File is not an archive")

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(str(drive_path / file.file_path), drive_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")

    async with _archive_semaphore:
        with zipfile.ZipFile(str(file_path), "r") as zf:
            # Look up entry by decoded name (handles Shift_JIS re-encoding)
            info = None
            for zi in zf.infolist():
                if _decode_zip_filename(zi) == path:
                    info = zi
                    break
            if info is None:
                raise HTTPException(
                    status_code=404, detail="Entry not found in archive"
                )

            # Reject symlink entries
            if info.external_attr != 0:
                mode = info.external_attr >> 16
                if mode != 0 and (mode & 0o170000) == 0o120000:
                    raise HTTPException(
                        status_code=400, detail="Symlink entries not supported"
                    )

            if info.file_size > _ARCHIVE_ENTRY_MAX_SIZE:
                raise HTTPException(
                    status_code=413,
                    detail="Entry exceeds maximum allowed size",
                )

            # Read with size limit to defend against ZIP bombs
            # (declared file_size can be spoofed in ZIP headers)
            with zf.open(info) as entry_fp:
                data = entry_fp.read(_ARCHIVE_ENTRY_MAX_SIZE + 1)
                if len(data) > _ARCHIVE_ENTRY_MAX_SIZE:
                    raise HTTPException(
                        status_code=413,
                        detail="Decompressed entry exceeds maximum allowed size",
                    )

    _, mime_type = classify(path)
    entry_filename = PurePosixPath(path).name
    disposition = "inline" if mime_type in _SAFE_INLINE_TYPES else "attachment"

    # RFC 6266: use filename* for non-ASCII names
    encoded_filename = quote(entry_filename, safe="")
    content_disp = f"{disposition}; filename*=UTF-8''{encoded_filename}"

    return Response(
        content=data,
        media_type=mime_type,
        headers={
            "Content-Length": str(len(data)),
            "Content-Disposition": content_disp,
            "X-Content-Type-Options": "nosniff",
        },
    )


@router.get("/{file_id}/subtitles/{index}")
async def get_subtitle(
    file_id: FileId,
    index: int,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)
    if file.file_type != "video" and file.mime_type != "application/vnd.litloft.loft+json":
        raise HTTPException(status_code=404, detail="Not a video file")

    drive_path = config.get_drive_path(file.drive)
    raw = detect_subtitles(file.file_path, drive_path)
    if index < 0 or index >= len(raw):
        raise HTTPException(status_code=404, detail="Subtitle not found")

    sub = raw[index]
    sub_path = _validate_path(str(drive_path / sub["path"]), drive_path)
    if not sub_path.exists():
        raise HTTPException(status_code=404, detail="Subtitle file not found on disk")

    _MAX_SUBTITLE_SIZE = 5 * 1024 * 1024  # 5MB
    if sub_path.stat().st_size > _MAX_SUBTITLE_SIZE:
        raise HTTPException(status_code=413, detail="Subtitle file too large")

    content = sub_path.read_text(encoding="utf-8-sig")

    if sub["format"] == "srt":
        content = convert_srt_to_vtt(content)

    return Response(
        content=content,
        media_type="text/vtt",
        headers={"Content-Type": "text/vtt; charset=utf-8"},
    )


@router.put("/{file_id}/rename", response_model=FileResponse)
async def rename_file_endpoint(
    file_id: FileId,
    body: FileRenameRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    _get_file_or_404(db, file_id, unlocked_groups)
    file = fileops.rename_file(db, file_id, body.new_filename)
    await event_hooks.emit("files.moved", {"file_ids": [file.id]})
    # Broadcast on the browser WebSocket alongside the addon webhook so
    # FolderTreePane / useFolderFiles subscribers refresh without a
    # manual reload. event_hooks dispatches to addon URLs only; the WS
    # layer is independent and reaches the browser directly.
    await ws_service.manager.broadcast(
        "files.moved", {"file_ids": [file.id]}, drive=file.drive
    )
    return _to_response(file)


@router.put("/{file_id}/move", response_model=FileResponse)
async def move_file_endpoint(
    file_id: FileId,
    body: FileMoveRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    _get_file_or_404(db, file_id, unlocked_groups)
    file = fileops.move_file(db, file_id, body.target_drive, body.target_folder_path)
    await event_hooks.emit("files.moved", {"file_ids": [file.id]})
    await ws_service.manager.broadcast(
        "files.moved", {"file_ids": [file.id]}, drive=file.drive
    )
    return _to_response(file)


@router.post("/{file_id}/copy", response_model=FileResponse)
async def copy_file_endpoint(
    file_id: FileId,
    body: FileCopyRequest,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    _get_file_or_404(db, file_id, unlocked_groups)
    new_file = fileops.copy_file(db, file_id, body.target_drive, body.target_folder_path)
    asyncio.create_task(
        event_hooks.emit("files.created", {"file_ids": [new_file.id]})
    )
    return _to_response(new_file)


@router.delete("/{file_id}")
async def delete_file_endpoint(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    _get_file_or_404(db, file_id, unlocked_groups)
    fileops.delete_file(db, file_id)
    asyncio.create_task(
        event_hooks.emit("files.deleted", {"file_ids": [file_id], "type": "soft_delete"})
    )
    return {"status": "deleted"}


@router.post("/{file_id}/restore", response_model=FileResponse)
async def restore_file_endpoint(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    _get_trashed_file_or_404(db, file_id, unlocked_groups)
    file = fileops.restore_file(db, file_id)
    asyncio.create_task(
        event_hooks.emit("files.restored", {"file_ids": [file_id]})
    )
    return _to_response(file)


def _compute_text_etag(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _inject_md_id(
    db: Session,
    file: File,
    body: bytes,
) -> tuple[bytes, str | None]:
    """Inject a frontmatter ``id:`` into ``body`` if missing.

    Returns ``(possibly_rewritten_body, new_md_id)``. ``new_md_id`` is
    ``None`` when no injection happened (malformed UTF-8/YAML, or the
    body already had a valid id — in which case ``new_md_id`` is still
    returned for projection so ``File.md_id`` stays in sync).

    Spec: docs/superpowers/specs/2026-05-12-markdown-link-three-forms.md §3.1.
    """
    try:
        body_str = body.decode("utf-8")
    except UnicodeDecodeError:
        return body, None

    parsed = parse_frontmatter(body_str)
    if not parsed.metadata:
        return body, None

    fm_id_raw = parsed.metadata.get("id")
    now = datetime.now(UTC)
    new_meta, new_id = ensure_id(
        parsed.metadata, existing_id=file.md_id, now=now
    )

    came_from_fm = isinstance(fm_id_raw, (str, int)) and str(fm_id_raw) == new_id
    came_from_db = (not came_from_fm) and file.md_id == new_id
    if not came_from_fm and not came_from_db:
        # Same-second collision insurance (spec §3.1): when the freshly
        # generated 14-digit id already exists as another file's md_id
        # in the same drive, append the millisecond component to extend
        # to 17 digits. 14-digit ids remain the common case.
        collision = (
            db.query(File.id)
            .filter(
                File.drive == file.drive,
                File.id != file.id,
                File.md_id == new_id,
            )
            .first()
        )
        if collision is not None:
            new_id = f"{new_id}{(now.microsecond // 1000):03d}"
            new_meta = {"id": new_id, **{k: v for k, v in new_meta.items() if k != "id"}}

    if str(parsed.metadata.get("id", "")) == new_id and "id" in parsed.metadata:
        return body, new_id

    new_body_str = compose_frontmatter(new_meta, parsed.body)
    return new_body_str.encode("utf-8"), new_id


def _strip_etag_quotes(value: str) -> str:
    value = value.strip()
    if value.startswith("W/"):
        value = value[2:].strip()
    if len(value) >= 2 and value[0] == '"' and value[-1] == '"':
        return value[1:-1]
    return value


@router.put("/{file_id}/content")
async def put_file_content(
    file_id: FileId,
    request: Request,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    """Overwrite a text file's content with optimistic-lock safety.

    Body: raw UTF-8 text, Content-Type: text/plain
    Required header: If-Match (ETag of current content)

    Responses:
    - 200 OK with new ETag on success
    - 412 if If-Match doesn't match current content's ETag
    - 413 if body exceeds size limit
    - 415 if target file's mime isn't in allowlist (text/markdown, text/plain)
    - 428 if If-Match is missing
    - 404 if file not found, trashed, or missing
    """
    if_match = request.headers.get("if-match")
    if not if_match:
        raise HTTPException(status_code=428, detail="If-Match header is required")

    body = await request.body()
    if len(body) > _TEXT_WRITE_MAX_BYTES:
        raise HTTPException(status_code=413, detail="Content exceeds size limit")

    file = _get_file_or_404(db, file_id, unlocked_groups)
    if (file.mime_type or "") not in _TEXT_WRITE_ALLOWED_MIMES:
        raise HTTPException(
            status_code=415,
            detail=f"Mime type not writable via this endpoint: {file.mime_type}",
        )

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(str(drive_path / file.file_path), drive_path)

    lock = _text_write_locks[file.id]
    async with lock:
        # Current ETag comes from actual file bytes (ETag is strong, content-hashed)
        try:
            current_bytes = file_path.read_bytes()
        except FileNotFoundError:
            raise HTTPException(status_code=404, detail="File not found on disk")
        current_etag = _compute_text_etag(current_bytes)

        if _strip_etag_quotes(if_match) != current_etag:
            raise HTTPException(status_code=412, detail="ETag mismatch")

        # Spec 2026-05-12-markdown-link-three-forms §3.1: inject frontmatter
        # ``id:`` for .md writes so wiki-link resolution has a stable handle.
        # Skipped silently on UnicodeDecodeError / malformed YAML — id is
        # never blocked on the write path.
        injected_md_id: str | None = None
        if _is_markdown_file(file):
            body, injected_md_id = _inject_md_id(db, file, body)
            if len(body) > _TEXT_WRITE_MAX_BYTES:
                raise HTTPException(
                    status_code=413, detail="Content exceeds size limit after id injection"
                )

        # Atomic write: tmp + os.replace. Temp file is on the same FS as target
        # so os.replace is atomic on POSIX.
        tmp_fd, tmp_name = tempfile.mkstemp(
            prefix=f".{file_path.name}.", suffix=".tmp", dir=str(file_path.parent)
        )
        try:
            with os.fdopen(tmp_fd, "wb") as f:
                f.write(body)
            os.replace(tmp_name, file_path)
        except Exception:
            try:
                os.unlink(tmp_name)
            except OSError:
                pass
            raise

        new_etag = _compute_text_etag(body)
        # Update File.file_size (ignore mtime — FS is authoritative for mtime)
        file.file_size = len(body)
        db.commit()

        # β canonical rule (spec 2026-04-24, Phase 11): for .md files,
        # the frontmatter's ``tags:`` is the source of truth for
        # ``File.tags``. Project synchronously so UI edits see the
        # effect without waiting for the knowledge scanner's hourly
        # pass and without a core → knowledge round-trip.
        #
        # Separate transaction so a projection failure (broken YAML,
        # invalid UTF-8, DB error on the tag write) cannot roll back
        # the content write that already landed on disk. Worst case
        # we log the error, the file bytes stay correct, and the
        # scanner reconciles ``File.tags`` within the hour.
        if _is_markdown_file(file):
            if injected_md_id is not None:
                try:
                    file.md_id = injected_md_id
                    db.commit()
                except Exception:
                    db.rollback()
                    logger.exception(
                        "put_content: md_id projection failed for %s", file_id
                    )
            try:
                parsed = parse_frontmatter(body.decode("utf-8"))
                tags = extract_valid_tags(parsed.metadata)
                replace_file_tags(db, file, tags)
                cleanup_orphan_tags(db)
                db.commit()
            except UnicodeDecodeError:
                db.rollback()
                logger.warning(
                    "put_content: %s is not valid UTF-8; skipping tag projection",
                    file_id,
                )
            except Exception:
                db.rollback()
                logger.exception(
                    "put_content: tag projection failed for %s", file_id
                )

            # Phase B (spec 2026-05-12 §3.6): project frontmatter
            # ``aliases:`` to ``File.md_aliases`` so the wiki-link
            # resolver can match alias-form targets. Isolated like the
            # tag projection — a parse / write failure must not roll
            # back the durable content write.
            try:
                parsed = parse_frontmatter(body.decode("utf-8"))
                aliases = extract_valid_aliases(parsed.metadata)
                file.md_aliases = json.dumps(aliases) if aliases else None
                db.commit()
            except UnicodeDecodeError:
                db.rollback()
                logger.warning(
                    "put_content: %s is not valid UTF-8; skipping aliases projection",
                    file_id,
                )
            except Exception:
                db.rollback()
                logger.exception(
                    "put_content: md_aliases projection failed for %s", file_id
                )

            # Sync loft:// + wiki-link relations → file_relations. Same
            # isolation pattern: a sync failure must not roll back the
            # content write. The resolver needs ``file.folder_path`` to
            # interpret ``./`` and ``../`` targets.
            try:
                content_str = body.decode("utf-8")
                _sync_md_file_relations(
                    db, file_id, file.drive, content_str, file.folder_path
                )
                db.commit()
            except UnicodeDecodeError:
                db.rollback()
                logger.warning(
                    "put_content: %s is not valid UTF-8; skipping link sync",
                    file_id,
                )
            except Exception:
                db.rollback()
                logger.exception(
                    "put_content: link sync failed for %s", file_id
                )

    asyncio.create_task(
        event_hooks.emit("files.updated", {"file_ids": [file_id]})
    )
    return Response(
        status_code=200,
        headers={"ETag": f'"{new_etag}"'},
    )


@router.get("/{file_id}/wiki-resolutions")
async def get_wiki_resolutions(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    """Return the resolver verdict for every ``[[X]]`` in a ``.md`` body.

    Spec ``2026-05-12-markdown-link-three-forms.md`` §3.8. Lets the
    renderer pick per-link styling (resolved / unresolved / ambiguous)
    without re-parsing the body in the browser.

    Response shape::

        {
          "resolutions": {
            "<target>": {"kind": "resolved", "file_id": "<id>"},
            "<target>": {"kind": "unresolved"},
            "<target>": {"kind": "ambiguous", "candidates": ["..."]}
          }
        }

    Errors: ``404`` for unknown / trashed / missing / inaccessible
    files; ``415`` when the file is not markdown. Drive-access gating
    flows through :func:`_get_file_or_404` so password-protected drives
    return ``404`` while locked.
    """
    file = _get_file_or_404(db, file_id, unlocked_groups)
    if not _is_markdown_file(file):
        raise HTTPException(status_code=415, detail="Not a markdown file")

    drive_path = config.get_drive_path(file.drive)
    file_path = _validate_path(str(drive_path / file.file_path), drive_path)
    try:
        content = file_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return {"resolutions": {}}

    extracted = extract_links(content)
    target_to_id, diagnostics = resolve_wiki_targets_with_map(
        db, file.drive, file.folder_path, extracted.wiki_targets
    )
    diagnostics_by_target = {d.target: d for d in diagnostics}

    # Bulk-load resolved files in one query so the renderer can show the
    # file's basename for id-form targets (``[[20260512143028]]``) instead
    # of the opaque id. Drive-scoped, active filter, .md predicate.
    resolved_ids = set(target_to_id.values())
    file_meta: dict[str, str] = {}
    if resolved_ids:
        rows = db.query(File.id, File.filename).filter(File.id.in_(resolved_ids)).all()
        file_meta = {r.id: r.filename for r in rows}

    resolutions: dict[str, dict] = {}
    # Preserve insertion order = body order for deterministic UI render.
    for target in extracted.wiki_targets:
        if target in resolutions:
            continue
        if target in target_to_id:
            fid = target_to_id[target]
            entry: dict = {"kind": "resolved", "file_id": fid}
            filename = file_meta.get(fid)
            if filename:
                # Basename without .md so the renderer can use it as
                # human-readable display text for id-form targets.
                basename = filename[:-3] if filename.lower().endswith(".md") else filename
                entry["filename"] = filename
                entry["basename"] = basename
            resolutions[target] = entry
        elif target in diagnostics_by_target:
            diag = diagnostics_by_target[target]
            entry = {"kind": diag.kind}
            if diag.kind == "ambiguous":
                entry["candidates"] = diag.candidates
            resolutions[target] = entry
    return {"resolutions": resolutions}


@router.delete("/{file_id}/purge")
async def purge_file_endpoint(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_trashed_or_missing_file_or_404(db, file_id, unlocked_groups)
    if file.missing_since is not None and file.deleted_at is None:
        fileops.purge_missing_file(db, file_id)
    else:
        fileops.purge_file(db, file_id)
    asyncio.create_task(
        event_hooks.emit("files.purged", {"file_ids": [file_id]})
    )
    return {"status": "purged"}


def _related_file_summary(file: File) -> RelatedFileSummary:
    return RelatedFileSummary(
        id=file.id,
        drive=file.drive,
        filename=file.filename,
        folder_path=file.folder_path,
        file_type=file.file_type,
        mime_type=file.mime_type,
        thumbnail_url=f"/api/files/{file.id}/thumbnail",
        has_thumbnail=file.thumbnail_path is not None,
        file_size=file.file_size,
        missing_since=file.missing_since,
        created_at=file.created_at,
        updated_at=file.updated_at,
    )


@router.get("/{file_id}/relations", response_model=FileRelationsResponse)
async def list_file_relations(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
    kind: Annotated[str | None, Query(max_length=32)] = None,
) -> FileRelationsResponse:
    """List files related to ``file_id`` via ``file_relations`` rows.

    The source file must be accessible to the caller (drive unlock).
    Results exclude related files that have been trashed; missing files
    are included so the UI can grey them out without dropping history.
    Related files live on the same drive as the source (spec R4), so
    access already covers both sides.
    """
    source = _get_file_or_404(db, file_id, unlocked_groups)

    q = db.query(FileRelation).filter(
        or_(
            FileRelation.file_id_a == file_id,
            FileRelation.file_id_b == file_id,
        )
    )
    if kind is not None:
        q = q.filter(FileRelation.kind == kind)
    q = q.order_by(FileRelation.created_at.desc())
    relations = q.all()
    if not relations:
        return FileRelationsResponse(relations=[])

    other_ids: list[str] = []
    relation_other: dict[int, str] = {}
    for rel in relations:
        other = rel.file_id_b if rel.file_id_a == file_id else rel.file_id_a
        relation_other[rel.id] = other
        other_ids.append(other)

    # Include missing files (missing_since set, deleted_at null) but drop
    # trashed ones. The UI wants a stable history with a greyed-out tile
    # rather than silent removal.
    other_files = (
        db.query(File)
        .filter(
            File.id.in_(other_ids),
            File.deleted_at.is_(None),
            File.drive == source.drive,
        )
        .all()
    )
    by_id = {f.id: f for f in other_files}

    items: list[FileRelationItem] = []
    for rel in relations:
        other_id = relation_other[rel.id]
        other = by_id.get(other_id)
        if other is None:
            continue
        items.append(
            FileRelationItem(
                relation_id=rel.id,
                kind=rel.kind,
                created_at=rel.created_at,
                created_by=rel.created_by,
                file=_related_file_summary(other),
            )
        )

    return FileRelationsResponse(relations=items)


@router.get("/{file_id}/exif", response_model=ExifResponse)
async def get_file_exif(
    file_id: FileId,
    db: Annotated[Session, Depends(get_db)],
    unlocked_groups: Annotated[list[str], Depends(get_unlocked_groups)],
):
    file = _get_file_or_404(db, file_id, unlocked_groups)
    if file.file_type != "image":
        raise HTTPException(status_code=404, detail="No EXIF data")
    exif = db.get(FileExif, file_id)
    if not exif:
        raise HTTPException(status_code=404, detail="No EXIF data")
    return ExifResponse.model_validate(exif)

